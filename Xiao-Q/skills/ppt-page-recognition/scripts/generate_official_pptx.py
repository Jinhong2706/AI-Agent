#!/usr/bin/env python3

import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# AI Generated Image Paths
IMG_EMS = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ems_icon_1775028844946.png"
IMG_IBMS = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ibms_icon_1775028858222.png"
IMG_AI_AGENT = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ai_agent_icon_1775028872392.png"
IMG_TIANMA = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/tianma_logo_placeholder_1775029023462.png"

def render_official_pptx(master_json_path: Path, clean_template_path: Path, output_path: Path):
    if not master_json_path.exists():
        print(f"Error: {master_json_path} not found.")
        return
    
    with open(master_json_path, 'r', encoding='utf-8') as f:
        master_data = json.load(f)
    
    # Load the "Clean Master" (0 slides, but all brand layouts intact)
    output_prs = Presentation(clean_template_path)
    
    # Layout Map (ID -> Name)
    # 0: 2_封面页
    # 1: 目录页
    # 2: 章节页
    # 3: 正文页
    # 5: 空白
    type_map = {
        "cover": 0,
        "toc": 1,
        "chapter": 2,
        "timeline": 5,    # Draw custom on Blank
        "logo-wall": 5,   # Draw custom on Blank
        "metric": 3,
        "fallback": 3
    }

    print(f"[*] Generating high-fidelity official PPTX for {len(master_data['slides'])} slides...")
    
    COLOR_DEEP_BLUE = RGBColor(0, 51, 153)
    COLOR_GOLD = RGBColor(197, 163, 100) # Brand gold

    for slide_data in master_data["slides"]:
        stype = slide_data.get("page_type", "fallback")
        # Special case labels for metrics or peers
        if stype in {"metrics", "peer-panels", "people-matrix", "awards-patents"}:
             stype = "metric"
             
        layout_idx = type_map.get(stype, 3) 
        dest_slide = output_prs.slides.add_slide(output_prs.slide_layouts[layout_idx])
        
        print(f"  [>] Slide {slide_data['slide_number']} (Type: {stype} -> Layout: {output_prs.slide_layouts[layout_idx].name})")
        
        # 1. Title
        title_text = slide_data.get("title", "")
        if dest_slide.shapes.title:
            dest_slide.shapes.title.text = title_text
            
        # 2. Content
        if stype == "timeline" and "structured_payload" in slide_data:
            payload = slide_data["structured_payload"]
            icons = [IMG_EMS, IMG_IBMS, IMG_AI_AGENT]
            periods = payload.get("periods", [])
            for i, period in enumerate(periods[:3]):
                left = Inches(0.5 + i * 3.1)
                top = Inches(2.0)
                # Icon
                if os.path.exists(icons[i]):
                     dest_slide.shapes.add_picture(icons[i], left + Inches(0.8), top, width=Inches(1.2))
                # Title
                txt_title = dest_slide.shapes.add_textbox(left, top + Inches(1.3), width=Inches(2.8), height=Inches(0.5))
                txt_title.text_frame.text = period.get("period", "")
                txt_title.text_frame.paragraphs[0].font.bold = True
                txt_title.text_frame.paragraphs[0].font.color.rgb = COLOR_DEEP_BLUE
                # Desc
                txt_desc = dest_slide.shapes.add_textbox(left, top + Inches(1.8), width=Inches(2.8), height=Inches(1.5))
                txt_desc.text_frame.word_wrap = True
                txt_desc.text_frame.text = period.get("headline", "")
                txt_desc.text_frame.paragraphs[0].font.size = Pt(12)

        elif stype == "logo-wall":
             if os.path.exists(IMG_TIANMA):
                  dest_slide.shapes.add_picture(IMG_TIANMA, Inches(3.5), Inches(2.2), width=Inches(3))
        
        else:
            # Fallback for Standard body
            for shape in dest_slide.placeholders:
                if shape.placeholder_format.idx in {13, 10}: # Content
                    body_lines = []
                    for item in slide_data.get("hierarchy", []):
                        body_lines.append(item.get("title", ""))
                        for sub in item.get("items", []):
                            body_lines.append(f"  • {sub}")
                    shape.text = "\n".join(body_lines)

    output_prs.save(output_path)
    print(f"\n[√] Official brand-style PPT generated: {output_path}")

if __name__ == "__main__":
    render_official_pptx(
        Path("output/master_deck.json"),
        Path("干净模板.pptx"),
        Path("项目成果_正式版.pptx")
    )
