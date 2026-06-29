#!/usr/bin/env python3

import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# AI Generated Image Paths
IMG_EMS = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ems_icon_1775028844946.png"
IMG_IBMS = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ibms_icon_1775028858222.png"
IMG_AI_AGENT = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/ai_agent_icon_1775028872392.png"
IMG_TIANMA = "/Users/mengdong/.gemini/antigravity/brain/50a8b7c9-2696-40e0-a821-3338c79dfe2d/tianma_logo_placeholder_1775029023462.png"

def clone_shape(shape, dest_slide):
    """Simple shape cloning function for common types."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # For pictures, we'd need to extract and re-add. 
            # For now, skip to avoid complex blob handling, or use a placeholder.
            pass
        elif shape.has_text_frame:
            # Create a text box at the same position
            new_shape = dest_slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text = shape.text
            # Copy basic font if possible
            if shape.text_frame.paragraphs:
                for i, p in enumerate(shape.text_frame.paragraphs):
                    if i < len(new_shape.text_frame.paragraphs):
                        new_p = new_shape.text_frame.paragraphs[i]
                    else:
                        new_p = new_shape.text_frame.add_paragraph()
                    new_p.text = p.text
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            # Clone simple auto-shapes (rectangles, etc.)
            new_shape = dest_slide.shapes.add_shape(
                shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height
            )
            # Match fill/line if needed (simplified)
            if hasattr(shape, "fill") and hasattr(new_shape, "fill"):
                 pass # Copy colors later
    except Exception as e:
        print(f"      [!] Error cloning shape {shape.name}: {e}")

def render_pptx(master_json_path: Path, template_path: Path, output_path: Path):
    if not master_json_path.exists():
        print(f"Error: {master_json_path} not found.")
        return
    
    with open(master_json_path, 'r', encoding='utf-8') as f:
        master_data = json.load(f)

    # We use a brand new, empty Presentation to avoid the zipfile/XML corruption 
    # bugs present in the original '活力结构模板.pptx' template.
    output_prs = Presentation()
    
    print(f"[*] Starting high-fidelity rendering for {len(master_data['slides'])} slides into a pristine blank presentation...")
    
    # Mapping based on standard default Blank python-pptx master template layouts
    # Layout indices in default template: 0=Title, 1=Title+Content, 6=Blank, etc.
    type_map = {
        "cover": 0,
        "toc": 1,
        "timeline": 6,   # Blank layout to manually draw
        "logo-wall": 6,  # Blank layout
        "peer-panels": 1,
        "people-matrix": 1,
        "awards-patents": 1,
        "metrics": 1,
        "chapter": 0,
        "fallback": 1,
    }

    for slide_data in master_data["slides"]:
        stype = slide_data.get("page_type", "fallback")
        template_idx = type_map.get(stype, 1) # Fallback to standard body
        
        # Use default layout from the blank presentation
        dest_layout = output_prs.slide_layouts[template_idx]
        dest_slide = output_prs.slides.add_slide(dest_layout)
        
        print(f"  [>] Processing Slide {slide_data['slide_number']} (Type: {stype} -> Default Layout {template_idx})")
        
        # 1. Fill Title
        title_text = slide_data.get("title", "")
        if dest_slide.shapes.title:
            dest_slide.shapes.title.text = title_text
            
        # 2. Fill Structured Payload
        payload = slide_data.get("structured_payload", {})
        
        if stype == "timeline" and "periods" in payload:
            print(f"    [*] Custom Rendering Slide 5 (Timeline) with generated icons...")
            icons = [IMG_EMS, IMG_IBMS, IMG_AI_AGENT]
            for i, period in enumerate(payload.get("periods", [])):
                if i >= 3: break # Limit to 3 columns
                left = Inches(1 + i * 3.3)
                top = Inches(2.2)
                # Add Icon
                if os.path.exists(icons[i]):
                     dest_slide.shapes.add_picture(icons[i], left, top, height=Inches(1.2))
                # Add Text
                txt_box = dest_slide.shapes.add_textbox(left, top + Inches(1.3), width=Inches(3), height=Inches(2))
                tf = txt_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"{period.get('period', '')}\n{period.get('headline', '')}"
                p.font.bold = True
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0, 51, 153) # Deep Blue
        
        if stype == "logo-wall" and os.path.exists(IMG_TIANMA):
             print(f"    [*] Inserting generated logo for Slide 6...")
             dest_slide.shapes.add_picture(IMG_TIANMA, Inches(4), Inches(2.5), height=Inches(1.8))
        
        # Fallback: Merge all text into the main content placeholder
        for shape in dest_slide.placeholders:
            if shape.placeholder_format.idx == 13 or shape.placeholder_format.idx == 10: # Content PH
                 body_lines = []
                 for item in slide_data.get("hierarchy", []):
                     body_lines.append(item.get("title", ""))
                     for sub in item.get("items", []):
                         body_lines.append(f"  • {sub}")
                 shape.text = "\n".join(body_lines)

    output_prs.save(output_path)
    print(f"\n[√] Final PPT generated: {output_path}")

if __name__ == "__main__":
    render_pptx(
        Path("output/master_deck.json"),
        Path("活力结构模板.pptx"),
        Path("白膜纯净版.pptx")
    )
